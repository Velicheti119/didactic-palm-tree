from pptx import Presentation
import os

def add_multiple_slides_from_template(template_path, slide_index, content_list, output_path):
    """
    Adds multiple slides to an existing PowerPoint file, based on a specific template slide.
    
    :param template_path: Path to the input PowerPoint template.
    :param slide_index: Index of the slide in the template to use as a base (0-based).
    :param content_list: List of content to add to each new slide.
    :param output_path: Path to save the modified PowerPoint file.
    """
    # Load the existing PowerPoint file
    prs = Presentation(template_path)
    
    # Get the slide layout to duplicate
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Invalid slide_index: {slide_index}. The template has {len(prs.slides)} slides.")
    
    template_slide = prs.slides[slide_index]

    # Add new slides using the content_list
    for i, content in enumerate(content_list):
        # Duplicate the selected slide
        slide_layout = template_slide.slide_layout
        new_slide = prs.slides.add_slide(slide_layout)
        
        # Add text to the slide's shapes (assumes placeholders are used in the template)
        for shape in new_slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:  # Modify title (first placeholder)
                shape.text = content.get("title", "")
            elif shape.is_placeholder and shape.placeholder_format.idx == 1:  # Modify content (second placeholder)
                shape.text = content.get("body", "")
        
        print(f"Slide {i + 1}: Added with title '{content.get('title', '')}'")

    # Save the modified presentation
    prs.save(output_path)
    print(f"Modified PowerPoint saved at '{output_path}'.")

# Example Usage
if __name__ == "__main__":
    template_file = "template.pptx"  # Path to your PPTX template file
    output_file = "output.pptx"  # Path to save the generated PPTX file
    content_data = [
        {"title": "Slide 1 Title", "body": "Slide 1 Content"},
        {"title": "Slide 2 Title", "body": "Slide 2 Content"},
        {"title": "Slide 3 Title", "body": "Slide 3 Content"}
    ]
    
    if os.path.exists(template_file):
        add_multiple_slides_from_template(template_file, 0, content_data, output_file)
    else:
        print(f"Template file '{template_file}' not found.")
